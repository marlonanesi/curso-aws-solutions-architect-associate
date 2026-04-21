from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor



# Definir caminho do arquivo PPTX final



# Criar apresentação Full HD
prs_retry = Presentation()
prs_retry.slide_width = Inches(13.33)
prs_retry.slide_height = Inches(7.5)
title_slide_layout = prs_retry.slide_layouts[0]
bullet_slide_layout = prs_retry.slide_layouts[1]

# Criar slides e adicionar conteúdo
for slide_data in slides_content:
    layout = title_slide_layout if not slide_data["content"] else bullet_slide_layout
    slide = prs_retry.slides.add_slide(layout)
    title = slide.shapes.title
    body = slide.placeholders[1] if slide_data["content"] else None

    title.text = slide_data["title"]
    if body:
        tf = body.text_frame
        for idx, line in enumerate(slide_data["content"]):
            if idx == 0:
                tf.text = line
            else:
                p = tf.add_paragraph()
                p.text = line

# Adicionar anotações visuais
for idx, slide_data in enumerate(slides_content):
    note = slide_data.get("note", "")
    if note:
        slide = prs_retry.slides[idx]
        left = Inches(0.5)
        top = Inches(6.5)
        width = Inches(12.5)
        height = Inches(1)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        p = tf.add_paragraph()
        p.text = note
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = RGBColor(100, 100, 100)

# Salvar apresentação
prs_retry.save(pptx_final_path)
print(f"Apresentação salva com sucesso: {pptx_final_path}")
